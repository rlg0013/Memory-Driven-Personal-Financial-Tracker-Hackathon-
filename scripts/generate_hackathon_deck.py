from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "backend"))

from memory_graph import get_memory_graph_data, get_monthly_summary, get_overspending_analysis  # noqa: E402


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "bg": RGBColor(7, 17, 13),
    "bg_soft": RGBColor(12, 24, 19),
    "surface": RGBColor(15, 29, 22),
    "surface_alt": RGBColor(24, 42, 33),
    "surface_light": RGBColor(242, 247, 240),
    "mint": RGBColor(184, 255, 118),
    "mint_deep": RGBColor(89, 160, 76),
    "cyan": RGBColor(141, 216, 255),
    "gold": RGBColor(241, 204, 106),
    "coral": RGBColor(255, 141, 118),
    "pink": RGBColor(255, 143, 197),
    "text": RGBColor(244, 246, 240),
    "text_soft": RGBColor(196, 210, 199),
    "text_dark": RGBColor(22, 38, 30),
    "green_dark": RGBColor(19, 51, 37),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(52, 83, 68),
}

FONT_TITLE = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Consolas"


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    logo_path = REPO_ROOT / "frontend" / "assets" / "logo1.jpeg"
    bg_path = REPO_ROOT / "frontend" / "assets" / "money_bg.jpeg"

    jan = get_monthly_summary("2026-01")
    feb = get_monthly_summary("2026-02")
    graph = get_memory_graph_data()
    overspending = get_overspending_analysis()

    total_income = jan["total_income"] + feb["total_income"]
    total_spent = jan["total_spent"] + feb["total_spent"]
    total_net = jan["net"] + feb["net"]
    top_categories = sorted(feb["by_category"].items(), key=lambda item: item[1], reverse=True)[:4]
    strongest_driver = overspending["drivers"][1]["detail"] if overspending.get("drivers") else overspending["summary"]

    add_cover_slide(prs, bg_path, logo_path)
    add_problem_slide(prs)
    add_solution_slide(prs, graph, total_income, total_spent, total_net)
    add_membrain_slide(prs, graph)
    add_features_slide(prs)
    add_live_data_slide(prs, jan, feb, graph, overspending, top_categories)
    add_advantages_slide(prs, strongest_driver)
    add_demo_slide(prs)

    output_dir = REPO_ROOT / "deliverables"
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "CatMoney_Hackathon_Pitch.pptx"
    prs.save(output_path)
    print(output_path)


def add_background(slide, color: RGBColor, image_path: Path | None = None, overlay_transparency: float = 0.08) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), 0, 0, width=SLIDE_W, height=SLIDE_H)
        overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = color
        overlay.fill.transparency = overlay_transparency
        overlay.line.fill.background()


def add_logo(slide, logo_path: Path, left: float, top: float, size: float = 0.7) -> None:
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(left), Inches(top), width=Inches(size), height=Inches(size))


def add_title(slide, title: str, subtitle: str | None = None, *, left: float = 0.7, top: float = 0.6, width: float = 7.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_TITLE
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLORS["text"]

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = FONT_BODY
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLORS["text_soft"]


def add_kicker(slide, text: str, *, left: float, top: float, width: float = 3.0, color: RGBColor | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = color or COLORS["mint"]


def add_paragraph(slide, text: str, *, left: float, top: float, width: float, height: float, size: int = 16, color: RGBColor | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.color.rgb = color or COLORS["text_soft"]


def add_card(slide, *, left: float, top: float, width: float, height: float, title: str, body: str, accent: RGBColor, dark: bool = True) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["surface"] if dark else COLORS["surface_light"]
    card.fill.transparency = 0.04 if dark else 0
    card.line.color.rgb = accent
    card.line.transparency = 0.4

    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    text_color = COLORS["text"] if dark else COLORS["text_dark"]
    body_color = COLORS["text_soft"] if dark else COLORS["green_dark"]

    box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.18), Inches(width - 0.3), Inches(height - 0.25))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_TITLE
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = text_color

    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    p2.line_spacing = 1.15
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = FONT_BODY
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = body_color


def add_bullets(slide, items: list[str], *, left: float, top: float, width: float, height: float, dark: bool = True, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        p.bullet = True
        r = p.add_run()
        r.text = item
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.color.rgb = COLORS["text_soft"] if dark else COLORS["green_dark"]


def add_stat_chip(slide, *, left: float, top: float, width: float, title: str, value: str, accent: RGBColor) -> None:
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.0))
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLORS["surface_alt"]
    chip.fill.transparency = 0.08
    chip.line.color.rgb = accent
    chip.line.transparency = 0.55

    box = slide.shapes.add_textbox(Inches(left + 0.14), Inches(top + 0.14), Inches(width - 0.24), Inches(0.72))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_BODY
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = COLORS["text_soft"]

    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = FONT_TITLE
    r2.font.size = Pt(18)
    r2.font.bold = True
    r2.font.color.rgb = COLORS["text"]


def add_cover_slide(prs: Presentation, bg_path: Path, logo_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg"], bg_path, overlay_transparency=0.18)

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.45), Inches(7.35), Inches(6.4))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(9, 21, 16)
    left_panel.fill.transparency = 0.18
    left_panel.line.color.rgb = COLORS["line"]
    left_panel.line.transparency = 0.6

    add_logo(slide, logo_path, 0.9, 0.82, 0.86)
    add_kicker(slide, "MEM-BRAIN HACKATHON DECK", left=1.86, top=0.92, width=2.8, color=COLORS["mint"])

    brand = slide.shapes.add_textbox(Inches(1.86), Inches(1.18), Inches(3.4), Inches(0.55))
    tf = brand.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "CatMoney"
    r.font.name = FONT_TITLE
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    title = slide.shapes.add_textbox(Inches(0.98), Inches(1.95), Inches(6.2), Inches(2.3))
    tf = title.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Memory-driven finance that explains behavior, not just balances."
    r.font.name = FONT_TITLE
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = COLORS["text"]

    sub = tf.add_paragraph()
    sub.space_before = Pt(14)
    rr = sub.add_run()
    rr.text = (
        "CatMoney turns fragmented SMS and transaction data into a living financial memory graph. "
        "It helps users understand overspending, see connected patterns, and recover savings goals with Mem-Brain at the core."
    )
    rr.font.name = FONT_BODY
    rr.font.size = Pt(13)
    rr.font.color.rgb = COLORS["text_soft"]

    chip_y = 5.25
    for idx, (label, color) in enumerate([
        ("Memory Graph", COLORS["cyan"]),
        ("Overspending AI", COLORS["mint"]),
        ("Goal Reasoning", COLORS["gold"]),
        ("Adaptive Tracker", COLORS["pink"]),
    ]):
        add_card(slide, left=0.98 + idx * 1.55, top=chip_y, width=1.4, height=0.8, title=label, body="", accent=color)

    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.15), Inches(4.45), Inches(5.2))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(245, 249, 242)
    right_panel.line.fill.background()

    stat_titles = [
        ("Mem-Brain role", "Core memory layer", COLORS["mint"]),
        ("Problem", "Fragmented money habits", COLORS["coral"]),
        ("Outcome", "Actionable financial memory", COLORS["cyan"]),
    ]
    for index, (label, value, accent) in enumerate(stat_titles):
        add_light_stat(slide, 8.55, 1.6 + index * 1.25, 3.75, label, value, accent)

    footer = slide.shapes.add_textbox(Inches(8.55), Inches(5.42), Inches(3.5), Inches(0.8))
    tf2 = footer.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "Built for live demo clarity, visual storytelling, and Mem-Brain-first technical depth."
    r2.font.name = FONT_BODY
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = COLORS["green_dark"]


def add_light_stat(slide, left: float, top: float, width: float, title: str, value: str, accent: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.95))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["surface_light"]
    card.line.color.rgb = accent
    card.line.transparency = 0.45

    box = slide.shapes.add_textbox(Inches(left + 0.16), Inches(top + 0.12), Inches(width - 0.3), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title.upper()
    r.font.name = FONT_BODY
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.color.rgb = COLORS["mint_deep"]

    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = FONT_TITLE
    r2.font.size = Pt(16)
    r2.font.bold = True
    r2.font.color.rgb = COLORS["green_dark"]


def add_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg_soft"])
    add_kicker(slide, "01  |  WHY THIS PRODUCT IS NEEDED", left=0.75, top=0.55, width=4.2)
    add_title(slide, "Most personal finance tools show transactions. They do not explain behavior.", left=0.75, top=0.9, width=7.2)

    add_paragraph(
        slide,
        "Users do not fail savings goals because they lack charts. They fail because spending patterns stay fragmented across small events, timing spikes, recurring merchants, and habits that ordinary apps do not remember or connect.",
        left=0.75,
        top=2.2,
        width=6.2,
        height=1.4,
        size=18,
    )

    add_card(slide, left=7.55, top=1.2, width=2.3, height=1.55, title="Dashboards stop at totals", body="They show where money went, but not why the same behavior repeats every month.", accent=COLORS["coral"])
    add_card(slide, left=10.05, top=1.2, width=2.3, height=1.55, title="Goals fail quietly", body="Static plans do not adapt when the user misses a checkpoint or saves extra.", accent=COLORS["gold"])
    add_card(slide, left=7.55, top=3.0, width=2.3, height=1.55, title="AI without memory is shallow", body="A plain chatbot answers from the prompt, not from linked financial context.", accent=COLORS["cyan"])
    add_card(slide, left=10.05, top=3.0, width=2.3, height=1.55, title="Patterns stay invisible", body="Recurring merchants, salary-week leaks, and goal conflicts rarely appear in one place.", accent=COLORS["pink"])

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.2), Inches(11.6), Inches(1.45))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["surface_alt"]
    band.line.color.rgb = COLORS["mint"]
    band.line.transparency = 0.7
    add_paragraph(
        slide,
        "What users actually need: a financial system that remembers context, surfaces relationships, and turns money history into decisions.",
        left=1.05,
        top=5.55,
        width=10.8,
        height=0.6,
        size=20,
        color=COLORS["text"],
    )


def add_solution_slide(prs: Presentation, graph: dict, total_income: float, total_spent: float, total_net: float) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg"])
    add_kicker(slide, "02  |  THE PRODUCT", left=0.75, top=0.55)
    add_title(slide, "CatMoney is a memory-driven finance assistant.", "It does not just log expenses; it stores financial behavior as linked memory and reasons over it.", left=0.75, top=0.9, width=7.0)

    add_card(slide, left=0.85, top=2.35, width=3.75, height=2.1, title="Remember", body="Every parsed transaction becomes a memory node with merchant, category, amount, date, month, week, and behavior metadata.", accent=COLORS["mint"])
    add_card(slide, left=4.8, top=2.35, width=3.75, height=2.1, title="Explain", body="The app uses linked memory to explain overspending, recurring categories, and conflict with savings goals.", accent=COLORS["cyan"])
    add_card(slide, left=8.75, top=2.35, width=3.75, height=2.1, title="Adapt", body="Goal checkpoints update dynamically when the user misses, partially completes, or exceeds a target.", accent=COLORS["gold"])

    add_stat_chip(slide, left=0.95, top=5.15, width=2.25, title="Income tracked", value=f"Rs {total_income:,.0f}", accent=COLORS["mint"])
    add_stat_chip(slide, left=3.4, top=5.15, width=2.25, title="Spend tracked", value=f"Rs {total_spent:,.0f}", accent=COLORS["coral"])
    add_stat_chip(slide, left=5.85, top=5.15, width=2.25, title="Net retained", value=f"Rs {total_net:,.0f}", accent=COLORS["cyan"])
    add_stat_chip(slide, left=8.3, top=5.15, width=2.0, title="Graph nodes", value=str(graph["summary"]["total_nodes"]), accent=COLORS["pink"])
    add_stat_chip(slide, left=10.5, top=5.15, width=2.0, title="Graph edges", value=str(graph["summary"]["total_edges"]), accent=COLORS["gold"])


def add_membrain_slide(prs: Presentation, graph: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg_soft"])
    add_kicker(slide, "03  |  MEM-BRAIN INTEGRATION", left=0.75, top=0.55, width=4.4)
    add_title(slide, "Mem-Brain is not a hidden backend dependency here. It is the product engine.", left=0.75, top=0.9, width=8.4)

    steps = [
        ("1. Ingest", "Mock bank SMS and transaction events are parsed into structured financial memories."),
        ("2. Enrich", "Each memory gets tags like merchant, category, amount, month, week, and transaction type."),
        ("3. Link", "Memories are connected to categories, merchants, time buckets, goals, and life events."),
        ("4. Retrieve", "Interpreted search and graph retrieval fetch context, not only keyword matches."),
        ("5. Reason", "The assistant and dashboards turn linked memories into insights, risks, and actions."),
    ]

    for idx, (title, body) in enumerate(steps):
        x = 0.9 + idx * 2.42
        add_card(slide, left=x, top=2.4, width=2.1, height=2.45, title=title, body=body, accent=[COLORS["mint"], COLORS["cyan"], COLORS["pink"], COLORS["gold"], COLORS["coral"]][idx])
        if idx < len(steps) - 1:
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 2.1), Inches(3.6), Inches(x + 2.35), Inches(3.6))
            connector.line.color.rgb = COLORS["mint"]
            connector.line.width = Pt(2.2)

    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(5.35), Inches(3.35), Inches(1.2))
    side.fill.solid()
    side.fill.fore_color.rgb = COLORS["surface_alt"]
    side.line.color.rgb = COLORS["cyan"]
    side.line.transparency = 0.55
    add_paragraph(
        slide,
        f"Current graph: {graph['summary']['transaction_nodes']} transaction memories, {graph['summary']['categories']} categories, {graph['summary']['merchants']} merchants, strongest relation = {graph['summary']['strongest_relation']}.",
        left=9.2,
        top=5.68,
        width=2.95,
        height=0.7,
        size=12,
        color=COLORS["text_soft"],
    )


def add_features_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg"])
    add_kicker(slide, "04  |  FEATURE SHOWCASE", left=0.75, top=0.55, width=3.8)
    add_title(slide, "The experience is built to make Mem-Brain visible in the UI.", left=0.75, top=0.9, width=8.0)

    features = [
        ("Overview Dashboard", "Live spending, income, category mix, and memory counts in a polished product shell.", COLORS["mint"]),
        ("Insights", "Monthly comparisons, top spend areas, and behavioral summaries from memory-backed analysis.", COLORS["cyan"]),
        ("Memory Graph", "A D3 force graph that visualizes how transactions connect to merchants, categories, and goals.", COLORS["pink"]),
        ("Overspending Explainer", "Answers why spending pressure appears, not just how much was spent.", COLORS["gold"]),
        ("Goal Reasoning", "Connects savings targets with real spending conflicts and action plans.", COLORS["coral"]),
        ("Adaptive Tracker", "Recalculates checkpoints when the user misses or exceeds a target so the goal can still be reached.", COLORS["cyan"]),
    ]

    positions = [(0.85, 2.1), (4.45, 2.1), (8.05, 2.1), (0.85, 4.35), (4.45, 4.35), (8.05, 4.35)]
    for (title, body, accent), (x, y) in zip(features, positions):
        add_card(slide, left=x, top=y, width=3.15, height=1.85, title=title, body=body, accent=accent)


def add_live_data_slide(prs: Presentation, jan: dict, feb: dict, graph: dict, overspending: dict, top_categories: list[tuple[str, float]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg_soft"])
    add_kicker(slide, "05  |  WHAT THE DEMO ALREADY PROVES", left=0.75, top=0.55, width=4.5)
    add_title(slide, "Realistic data makes the product story dynamic, not staged.", left=0.75, top=0.9, width=7.9)

    add_stat_chip(slide, left=0.85, top=1.95, width=2.55, title="Jan net balance", value=f"Rs {jan['net']:,.0f}", accent=COLORS["mint"])
    add_stat_chip(slide, left=3.6, top=1.95, width=2.55, title="Feb net balance", value=f"Rs {feb['net']:,.0f}", accent=COLORS["cyan"])
    add_stat_chip(slide, left=6.35, top=1.95, width=2.55, title="Top cluster", value=graph["summary"]["top_cluster"].title(), accent=COLORS["pink"])
    add_stat_chip(slide, left=9.1, top=1.95, width=2.55, title="Top merchant", value=graph["summary"]["top_recurring_merchant"].title(), accent=COLORS["gold"])

    add_card(slide, left=0.85, top=3.25, width=5.35, height=2.7, title="Behavior insight", body=overspending["summary"], accent=COLORS["mint"])
    add_card(
        slide,
        left=6.45,
        top=3.25,
        width=5.2,
        height=2.7,
        title="Why this matters in demo",
        body=(
            "The deck is backed by 48 natural transaction memories across two positive-balance months. "
            "That gives the graphs, comparisons, overspending logic, and goal reasoning enough variation to feel real."
        ),
        accent=COLORS["cyan"],
    )

    add_bar_panel(slide, top_categories, left=9.0, top=4.0, width=2.2)


def add_bar_panel(slide, categories: list[tuple[str, float]], *, left: float, top: float, width: float) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.55))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["surface"]
    panel.line.color.rgb = COLORS["line"]
    panel.line.transparency = 0.5

    max_amount = max((amount for _, amount in categories), default=1)
    title = slide.shapes.add_textbox(Inches(left + 0.14), Inches(top + 0.08), Inches(width - 0.2), Inches(0.25))
    ttf = title.text_frame
    p = ttf.paragraphs[0]
    r = p.add_run()
    r.text = "FEB TOP CATEGORIES"
    r.font.name = FONT_BODY
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.color.rgb = COLORS["text_soft"]

    accents = [COLORS["mint"], COLORS["cyan"], COLORS["pink"], COLORS["gold"]]
    for idx, (name, amount) in enumerate(categories):
        y = top + 0.42 + idx * 0.28
        label = slide.shapes.add_textbox(Inches(left + 0.14), Inches(y), Inches(0.95), Inches(0.18))
        tf = label.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = name.replace("_", " ").title()
        r.font.name = FONT_BODY
        r.font.size = Pt(8.5)
        r.font.color.rgb = COLORS["text_soft"]

        base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.98), Inches(y + 0.03), Inches(0.92), Inches(0.11))
        base.fill.solid()
        base.fill.fore_color.rgb = COLORS["surface_alt"]
        base.line.fill.background()

        fill_w = 0.92 * (amount / max_amount)
        fill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.98), Inches(y + 0.03), Inches(fill_w), Inches(0.11))
        fill.fill.solid()
        fill.fill.fore_color.rgb = accents[idx % len(accents)]
        fill.line.fill.background()


def add_advantages_slide(prs: Presentation, strongest_driver: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg"])
    add_kicker(slide, "06  |  WHY CATMONEY IS STRONG FOR A HACKATHON", left=0.75, top=0.55, width=5.0)
    add_title(slide, "The advantage is not just AI. The advantage is memory-backed reasoning.", left=0.75, top=0.9, width=8.4)

    compare_left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.0), Inches(5.6), Inches(4.6))
    compare_left.fill.solid()
    compare_left.fill.fore_color.rgb = COLORS["surface"]
    compare_left.line.color.rgb = COLORS["coral"]
    compare_left.line.transparency = 0.55
    add_title(slide, "Typical finance app", left=1.1, top=2.25, width=2.8)
    add_bullets(
        slide,
        [
            "Shows totals, filters, and category charts.",
            "Treats transactions as rows, not relationships.",
            "Savings plans stay static even after a miss.",
            "Chatbots answer from the prompt, not long-term context.",
        ],
        left=1.1,
        top=2.9,
        width=4.8,
        height=3.0,
        size=16,
    )

    compare_right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(2.0), Inches(5.6), Inches(4.6))
    compare_right.fill.solid()
    compare_right.fill.fore_color.rgb = COLORS["surface_alt"]
    compare_right.line.color.rgb = COLORS["mint"]
    compare_right.line.transparency = 0.4
    add_title(slide, "CatMoney", left=7.1, top=2.25, width=2.3)
    add_bullets(
        slide,
        [
            "Stores each financial event as linked memory.",
            "Explains why overspending happens and what to cut first.",
            "Revises checkpoint plans to still reach the end goal.",
            "Uses Mem-Brain search, graph operations, and memory linking visibly in the UI.",
        ],
        left=7.1,
        top=2.9,
        width=4.8,
        height=3.0,
        size=16,
    )

    add_paragraph(
        slide,
        f"Current strongest behavioral signal: {strongest_driver}",
        left=7.1,
        top=5.95,
        width=4.7,
        height=0.45,
        size=11,
        color=COLORS["text_soft"],
    )


def add_demo_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["bg_soft"])
    add_kicker(slide, "07  |  PRESENTATION FLOW", left=0.75, top=0.55, width=3.8)
    add_title(slide, "A clean 2-minute demo story for judges", left=0.75, top=0.9, width=6.2)

    steps = [
        "Start with the problem: current apps track money, but they do not remember financial behavior.",
        "Open the Memory Graph to show that transactions, merchants, categories, and goals are linked memories.",
        "Use the Overspending tab to explain why spending pressure happens in plain language.",
        "Use Goal Reasoning and the adaptive tracker to show how missed checkpoints are recalculated.",
        "Close with reliability: even when upstream AI is rate-limited, CatMoney still answers from stored memory.",
    ]
    add_bullets(slide, steps, left=0.95, top=2.0, width=6.1, height=4.7, dark=True, size=17)

    highlight = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.35), Inches(1.7), Inches(5.0), Inches(4.9))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = COLORS["surface_light"]
    highlight.line.color.rgb = COLORS["mint_deep"]
    highlight.line.transparency = 0.45

    box = slide.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(4.2), Inches(4.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Final message to judges"
    r.font.name = FONT_TITLE
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLORS["green_dark"]

    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    p2.line_spacing = 1.25
    r2 = p2.add_run()
    r2.text = (
        "CatMoney is not another budgeting dashboard.\n\n"
        "It is a memory-driven finance assistant that uses Mem-Brain to connect spending, patterns, goals, and decisions in one coherent product."
    )
    r2.font.name = FONT_BODY
    r2.font.size = Pt(15)
    r2.font.color.rgb = COLORS["green_dark"]

    p3 = tf.add_paragraph()
    p3.space_before = Pt(12)
    r3 = p3.add_run()
    r3.text = "That is the story to land."
    r3.font.name = FONT_MONO
    r3.font.size = Pt(15)
    r3.font.bold = True
    r3.font.color.rgb = COLORS["mint_deep"]


if __name__ == "__main__":
    main()
